#!/usr/bin/env python3
"""
Add a "Why Me" slide to the existing PPT.
Inserts as Slide 10 (before Research Insights), pushes others down.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree

# ── Color Palette (same as original) ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE  = RGBColor(0x2E, 0x6B, 0x9E)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x20, 0x20, 0x20)
GRAY      = RGBColor(0x60, 0x60, 0x60)
ACCENT_ORANGE = RGBColor(0xE8, 0x7C, 0x2A)
ACCENT_GREEN  = RGBColor(0x2D, 0x9C, 0x6F)

prs = Presentation('Research_Experience_Presentation.pptx')
W = prs.slide_width


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_section_header(slide, text, sub_text=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
                 text, font_size=28, color=WHITE, bold=True, font_name='Calibri')
    if sub_text:
        add_text_box(slide, Inches(0.6), Inches(0.7), Inches(10), Inches(0.4),
                     sub_text, font_size=14, color=RGBColor(0xA0, 0xC8, 0xE8),
                     font_name='Calibri')
    add_rect(slide, Inches(0), Inches(1.1), W, Inches(0.04), ACCENT_ORANGE)


# ══════════════════════════════════════════════════════════════
# Create the "Why Me" slide (appended at end, we'll reorder later)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# White background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_section_header(slide, 'Why Me — Competitive Strengths',
                   'What I bring to your Robotics M.S. program')

# ── Left Column: Core Strengths ──
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             'Core Strengths', font_size=20, color=DARK_BLUE, bold=True)

strengths = [
    ('Cross-Disciplinary Full-Stack',
     'EME background + AI/VLM + ROS2 software + embedded deployment (Jetson Orin Nano). '
     'Can bridge the gap between mechanical design and intelligent software — a rare combination.'),
    ('Real Hardware, Not Just Simulation',
     'Built and debugged a complete autonomous navigation system on a physical robot. '
     'Solved real-world issues: 3 depth camera swaps, TF2 coordinate bugs, PID oscillation, '
     'QoS compatibility — problems that only surface on real hardware.'),
    ('Cutting-Edge Technology Adoption',
     'Implemented VLM Function Calling Agent architecture on a real robot — a technique that '
     'appeared in academic papers only in 2024-2025. Demonstrates ability to rapidly learn '
     'and apply frontier research.'),
    ('Systems Thinking & Modular Design',
     'Designed a 4-layer architecture from scratch. Created a Skill YAML framework enabling '
     'zero-code behavioral extension. Shows I think beyond immediate problems toward scalable solutions.'),
]

y = Inches(2.1)
for title, desc in strengths:
    # Accent bar
    add_rect(slide, Inches(0.6), y, Inches(0.08), Inches(0.85), MED_BLUE)
    add_text_box(slide, Inches(0.9), y, Inches(5.3), Inches(0.35),
                 title, font_size=14, color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(0.9), y + Inches(0.3), Inches(5.3), Inches(0.55),
                 desc, font_size=11, color=GRAY)
    y += Inches(1.05)

# ── Right Column: Technical Skills + Alignment ──
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
             'Technical Skills', font_size=20, color=DARK_BLUE, bold=True)

# Skills grid
skills = [
    ('Robotics', 'ROS2 · PID Control · Pure Pursuit\nTF2 · Sensor Fusion'),
    ('Programming', 'Python · Ubuntu/Linux\nGit · API Integration'),
    ('AI / ML', 'VLM (Gemini) · Function Calling\nPrompt Engineering · Agent Design'),
    ('Engineering', 'SolidWorks · CAD · Unity\nJetson Orin Nano · Embedded'),
]

y_skill = Inches(2.1)
for cat, items in skills:
    add_rect(slide, Inches(7.0), y_skill, Inches(2.3), Inches(0.85),
             LIGHT_BLUE)
    add_text_box(slide, Inches(7.15), y_skill + Inches(0.05), Inches(2.1), Inches(0.3),
                 cat, font_size=12, color=DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.15), y_skill + Inches(0.3), Inches(2.1), Inches(0.5),
                 items, font_size=10, color=GRAY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    y_skill += Inches(0.95)

# Project count highlight
add_rect(slide, Inches(9.6), Inches(2.1), Inches(3.0), Inches(3.8),
         RGBColor(0xFE, 0xF3, 0xE2))

add_text_box(slide, Inches(9.6), Inches(2.2), Inches(3.0), Inches(0.4),
             'By the Numbers', font_size=14, color=DARK_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

numbers = [
    ('4', 'Research Projects'),
    ('3+', 'Months of Intensive\nHardware Debugging'),
    ('10+', 'Major System Iterations'),
    ('4', 'Layer Architecture\nDesigned from Scratch'),
    ('7', 'Atomic Tool Functions\nfor Agent System'),
]

y_num = Inches(2.7)
for num, label in numbers:
    add_text_box(slide, Inches(9.8), y_num, Inches(0.7), Inches(0.4),
                 num, font_size=20, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(10.55), y_num + Inches(0.02), Inches(1.9), Inches(0.4),
                 label, font_size=10, color=GRAY, line_spacing=1.3)
    y_num += Inches(0.55)

# ── Bottom: Alignment Statement ──
add_rect(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.9),
         RGBColor(0xE8, 0xF0, 0xF8))
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.5), Inches(0.7),
             'Alignment with Robotics M.S.:  My hands-on experience in VLM-driven autonomous systems, '
             'real hardware deployment, and modular software architecture directly prepares me to contribute '
             'to cutting-edge robotics research from day one. I am eager to deepen my expertise in embodied '
             'AI and multi-robot systems at the graduate level.',
             font_size=13, color=DARK_BLUE, bold=False)


# ══════════════════════════════════════════════════════════════
# Move the new slide to position 10 (index 9), before "Research Insights"
# The new slide was appended as the last slide (index 11, i.e., slide 12)
# We want it at index 9 (slide 10)
# ══════════════════════════════════════════════════════════════

# python-pptx doesn't have a native move_slide function,
# so we manipulate the XML directly
slide_list = prs.slides._sldIdLst
slide_elements = list(slide_list)
# The new slide is the last element
new_slide_element = slide_elements[-1]
# We want to insert before index 9 (0-based), which is the 10th slide
# Current order: 0-Title, 1-AboutMe, 2-Early1, 3-Early2, 4-FYP, 5-Arch,
#                6-Modes, 7-Contributions, 8-Experiment, 9-Insights, 10-ThankYou, 11-NEW
# Target: insert NEW at index 9, pushing Insights to 10 and ThankYou to 11
target_element = slide_elements[9]  # "Research Insights" slide
slide_list.remove(new_slide_element)
slide_list.insert(slide_list.index(target_element), new_slide_element)

# Save
output_path = '/home/kuko/humble_ws/Research_Experience_Presentation.pptx'
prs.save(output_path)
print(f'✅ "Why Me" slide added and saved to: {output_path}')
print(f'   Total slides: {len(prs.slides)}')
print(f'   New slide position: Slide 10 (before Research Insights)')
