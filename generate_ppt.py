#!/usr/bin/env python3
"""
Generate a research presentation PPT for scholarship application.
Theme: VLM-Driven Navigation for Heterogeneous Multi-Robot Systems
Style: Clean academic presentation (~10 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (professional academic blue theme) ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)    # title bg / accent
MED_BLUE  = RGBColor(0x2E, 0x6B, 0x9E)     # section headers
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)    # light bg accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x20, 0x20, 0x20)
GRAY      = RGBColor(0x60, 0x60, 0x60)
ACCENT_ORANGE = RGBColor(0xE8, 0x7C, 0x2A) # highlight
ACCENT_GREEN  = RGBColor(0x2D, 0x9C, 0x6F)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ══════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════

def add_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, items,
                              font_size=16, color=BLACK, bullet_color=MED_BLUE,
                              font_name='Calibri', line_spacing=1.5):
    """Add bulleted list content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check for sub-item (starts with '  -')
        is_sub = item.startswith('  -')
        text = item.lstrip(' -').strip()

        if is_sub:
            p.text = '    ▸  ' + text
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = GRAY
        else:
            p.text = '●  ' + text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color

        p.font.name = font_name
        p.space_after = Pt(4)
        p.line_spacing = Pt(int(font_size * line_spacing))

    return txBox


def add_section_header(slide, text, sub_text=None):
    """Add a consistent section header bar at top of slide."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
                 text, font_size=28, color=WHITE, bold=True, font_name='Calibri')
    if sub_text:
        add_text_box(slide, Inches(0.6), Inches(0.7), Inches(10), Inches(0.4),
                     sub_text, font_size=14, color=RGBColor(0xA0, 0xC8, 0xE8),
                     font_name='Calibri')
    # Bottom accent line
    add_rect(slide, Inches(0), Inches(1.1), W, Inches(0.04), ACCENT_ORANGE)


def add_page_number(slide, num, total):
    """Add page number at bottom right."""
    add_text_box(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                 f'{num}/{total}', font_size=11, color=GRAY,
                 alignment=PP_ALIGN.RIGHT, font_name='Calibri')


TOTAL_SLIDES = 10


# ══════════════════════════════════════════════════════════════
# Slide 1: Title Slide
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Decorative top bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.15), ACCENT_ORANGE)

# Title
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.5),
             'Research Experience Overview',
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Calibri')

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.8),
             'VLM-Driven Visual Navigation for Heterogeneous Multi-Robot Systems',
             font_size=22, color=RGBColor(0xA0, 0xC8, 0xE8), bold=False,
             alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Separator line
add_rect(slide, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.03), ACCENT_ORANGE)

# Placeholder for personal info
add_text_box(slide, Inches(2.0), Inches(4.7), Inches(9.3), Inches(2.0),
             '[Your Name]\n'
             '[University / Department]\n'
             '[Email Address]\n'
             '[Date]',
             font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
             alignment=PP_ALIGN.CENTER, font_name='Calibri', line_spacing=1.8)

add_page_number(slide, 1, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 2: About Me / Background (placeholder for user)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'About Me')

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             'Education', font_size=22, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.0), [
    '[Degree], [University], [Year]',
    '  - Major: [Your Major]',
    '  - GPA: [X.XX / 4.0]',
    '[Other relevant education or exchange experience]',
], font_size=16)

add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.5),
             'Research Interests', font_size=22, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(2.0), [
    'Multimodal AI × Robotics',
    'Vision-Language Models for Embodied Intelligence',
    'Autonomous Navigation & Semantic Mapping',
    'Multi-Robot Collaboration Systems',
], font_size=16)

# Right side: placeholder for photo / additional info
add_rect(slide, Inches(7.5), Inches(1.5), Inches(4.5), Inches(5.2),
         LIGHT_BLUE)
add_text_box(slide, Inches(7.5), Inches(3.5), Inches(4.5), Inches(1.0),
             '[Photo / Skills / Awards\nPlaceholder]',
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 2, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 3: Early Research Experience (placeholder for user)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Early Research Experience',
                   '[Fill in your earlier research projects before the FYP]')

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
             'Project 1: [Title]', font_size=22, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.0), [
    'Objective: [Describe the research goal]',
    'Method: [Key approach / algorithm / framework]',
    'Outcome: [Key result or contribution]',
    '  - [Relevant skill gained: e.g., Python, MATLAB, data analysis]',
], font_size=16)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
             'Project 2: [Title] (optional)', font_size=22, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.0), [
    'Objective: [Describe the research goal]',
    'Method: [Key approach]',
    'Outcome: [Key result]',
], font_size=16)

add_page_number(slide, 3, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 4: FYP Overview — Motivation & Problem Statement
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Final Year Project: Motivation & Problem',
                   'VLM-Driven Navigation for Heterogeneous Multi-Robot Systems')

# Left column: Problem
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             'Problem Statement', font_size=20, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), [
    'Traditional robot navigation relies on pre-built maps (SLAM) and hand-coded rules — brittle in unknown environments',
    'Most VLM-robotics studies remain in simulation; real hardware deployment is rare',
    'Heterogeneous multi-robot collaboration with LLMs is largely unexplored',
    'Need: a map-free, vision-driven navigation system validated on real hardware',
], font_size=15)

# Right column: Approach
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
             'Our Approach', font_size=20, color=ACCENT_GREEN, bold=True)

add_bullet_slide_content(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5), [
    'Replace SLAM with Gemini 2.5 Flash VLM as the cognitive planner',
    'Two-stage architecture: single-turn pipeline → multi-turn Function Calling Agent',
    'Skill-based YAML framework for task-specific prompts',
    'Deployed on real robot (Jetson Orin Nano + depth camera)',
    'Leader-Follower heterogeneous formation',
], font_size=15, color=BLACK)

# Bottom highlight box
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), LIGHT_BLUE)
add_text_box(slide, Inches(1.2), Inches(5.9), Inches(11.0), Inches(0.8),
             'Key Insight:  A well-crafted prompt + function-calling tools can replace complex hand-tuned '
             'planning algorithms — achieving real-world indoor navigation without any pre-built map.',
             font_size=15, color=DARK_BLUE, bold=False, font_name='Calibri')

add_page_number(slide, 4, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 5: System Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'System Architecture',
                   'Four-layer modular design on ROS2 Humble')

# Architecture layers - visual representation
layers = [
    ('Cognitive Planning Layer', 'Gemini 2.5 Flash API  |  Skill YAML Profiles  |  Function Calling Agent Loop', DARK_BLUE, WHITE),
    ('Perception Layer', 'RGB/Depth Processing  |  TF2 Coordinate Transform  |  Semantic Map Cache', MED_BLUE, WHITE),
    ('Control Execution Layer', 'Pure Pursuit + PID Tracker  |  Reactive Collision Avoidance (Depth ROI)', RGBColor(0x45, 0x8B, 0x74), WHITE),
    ('Hardware Abstraction Layer', 'Mecanum Chassis + STM32  |  Orbbec Gemini 336 Depth Camera  |  IMU  |  Jetson Orin Nano', RGBColor(0x6B, 0x6B, 0x6B), WHITE),
]

y_start = Inches(1.5)
layer_h = Inches(1.15)
gap = Inches(0.15)
for i, (name, desc, bg_color, txt_color) in enumerate(layers):
    y = y_start + i * (layer_h + gap)
    add_rect(slide, Inches(0.8), y, Inches(7.5), layer_h, bg_color)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(3.0), Inches(0.4),
                 name, font_size=16, color=txt_color, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.55), Inches(7.0), Inches(0.5),
                 desc, font_size=13, color=RGBColor(0xE0, 0xE0, 0xE0))

    # Arrow between layers
    if i < len(layers) - 1:
        add_text_box(slide, Inches(4.2), y + layer_h - Inches(0.05), Inches(0.5), Inches(0.35),
                     '▼', font_size=16, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# Right side: data flow
add_text_box(slide, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.4),
             'End-to-End Data Flow', font_size=17, color=DARK_BLUE, bold=True)

flow_items = [
    '1. RGB frame → Gemini API',
    '2. Gemini returns pixel coords / odom waypoints',
    '3. Depth reprojection → 3D path (single-turn)\n'
    '    or direct odom coords (agent mode)',
    '4. PID tracker → /cmd_vel',
    '5. Obstacle detected → re-trigger Gemini',
    '6. Semantic labels stored in memory',
]
add_bullet_slide_content(slide, Inches(9.0), Inches(2.1), Inches(4.0), Inches(4.5),
                          flow_items, font_size=13, color=BLACK)

add_page_number(slide, 5, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 6: Two Operating Modes
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Two Operating Modes',
                   'Single-Turn Pipeline vs. Multi-Turn Function Calling Agent')

# Left: Single-Turn
add_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.5),
             'Mode A: Single-Turn Pipeline', font_size=19, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.3), Inches(5.4), Inches(4.0), [
    'track_path triggers → capture RGB image',
    'Send image to Gemini → returns pixel coordinates',
    'image_conversion: depth reprojection + TF2 → 3D odom path',
    'PID tracker follows path → re-trigger on completion',
    '  - Pros: Simple, low API cost per call',
    '  - Cons: Passive — Gemini cannot query robot state',
    '  - Cons: Sensitive to depth camera data quality',
], font_size=14)

# Right: Agent Mode
add_rect(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
         RGBColor(0xE8, 0xF5, 0xE0))
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.5),
             'Mode B: Function Calling Agent', font_size=19, color=ACCENT_GREEN, bold=True)

add_bullet_slide_content(slide, Inches(7.2), Inches(2.3), Inches(5.4), Inches(4.0), [
    'Gemini operates as an interactive agent',
    'Actively calls tools: get_robot_pose(), get_front_image(), publish_path()',
    'Computes odom waypoints via trigonometry — no depth reprojection needed',
    'Multi-turn loop until task complete or max turns reached',
    '  - Pros: Spatially grounded — knows exact pose',
    '  - Pros: Generates semantic labels during navigation',
    '  - Cons: Higher API token consumption (~3-5× per cycle)',
], font_size=14)

add_page_number(slide, 6, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 7: Key Technical Contributions
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Key Technical Contributions')

contributions = [
    ('Skill-Based YAML Framework',
     'System prompts defined in hot-swappable YAML files. Adding new robot behaviors '
     '(exploration, inspection, navigation) requires zero code changes — just a new YAML file.',
     '🔧'),
    ('Prompt Engineering for Perspective Compensation',
     'Camera mounted at only 155mm height causes severe perspective compression. '
     'Engineered prompts instruct Gemini to place far waypoints at image y≈500-550, '
     'achieving 1.5-2m per step (vs. <0.5m with naive prompt).',
     '📐'),
    ('Function Calling Agent Architecture',
     '7 atomic tool functions (get_pose, get_image, publish_path, label_region, etc.) '
     'enable Gemini to actively query sensors and make spatially grounded decisions. '
     'Bypasses depth camera reprojection entirely.',
     '🤖'),
    ('Real Hardware Deployment',
     'Complete system validated on physical robot: Jetson Orin Nano + Orbbec depth camera + '
     'Mecanum chassis. Iterative debugging of depth camera issues, TF2 transforms, '
     'PID tuning, and QoS compatibility.',
     '⚙️'),
]

y = Inches(1.5)
for emoji, (title, desc, _) in zip(['🔧', '📐', '🤖', '⚙️'], contributions):
    # Colored left bar
    add_rect(slide, Inches(0.6), y, Inches(0.1), Inches(1.1), MED_BLUE)

    add_text_box(slide, Inches(1.0), y, Inches(11.5), Inches(0.4),
                 title, font_size=17, color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.4), Inches(11.5), Inches(0.7),
                 desc, font_size=13, color=GRAY)
    y += Inches(1.35)

add_page_number(slide, 7, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 8: Experimental Validation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Experimental Validation',
                   '3-Group Comparative Study: Naive Prompt vs. Optimized Prompt vs. Agent Mode')

# Experiment setup
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             'Experiment Setup', font_size=18, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5), [
    'Scenario: 8m corridor with obstacle, target = exit door',
    '5 trials per group, same starting pose',
    'Group A: Naive prompt (single-turn)',
    'Group B: Optimized prompt (single-turn)',
    'Group C: Function Calling Agent (multi-turn)',
], font_size=14)

# Results table header
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.4),
             'Key Metrics (to be filled after experiments)', font_size=18, color=DARK_BLUE, bold=True)

# Table-like structure
table_data = [
    ('Metric', 'Group A', 'Group B', 'Group C'),
    ('Success Rate', '[___]', '[___]', '[___]'),
    ('Avg. Time (s)', '[___]', '[___]', '[___]'),
    ('Gemini Interactions', '[___]', '[___]', '[___]'),
    ('Angular Vel. Var.', '[___]', '[___]', '[___]'),
    ('Semantic Labels', 'N/A', 'N/A', '[___]'),
]

y = Inches(5.0)
for i, row in enumerate(table_data):
    for j, cell in enumerate(row):
        x = Inches(0.8 + j * 1.6)
        bg = DARK_BLUE if i == 0 else (LIGHT_BLUE if i % 2 == 0 else WHITE)
        tc = WHITE if i == 0 else BLACK
        bld = (i == 0)
        fs = 12 if i == 0 else 12
        add_rect(slide, x, y, Inches(1.55), Inches(0.35), bg)
        add_text_box(slide, x + Inches(0.05), y + Inches(0.02), Inches(1.45), Inches(0.3),
                     cell, font_size=fs, color=tc, bold=bld, alignment=PP_ALIGN.CENTER)
    y += Inches(0.35)

# Right side: figure placeholder
add_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2),
         LIGHT_BLUE)
add_text_box(slide, Inches(7.2), Inches(3.3), Inches(5.5), Inches(1.5),
             '[Figure: Trajectory Comparison Plot\n'
             'Group A (red) / Group B (blue) / Group C (green)\n'
             'on odom coordinate axes]',
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 8, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 9: Research Insights & Broader Impact
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 'Research Insights & Future Directions')

# Left: Insights
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             'Key Takeaways', font_size=20, color=DARK_BLUE, bold=True)

add_bullet_slide_content(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), [
    'Cloud VLMs can serve as practical path planners for real robots — not just simulation',
    'Prompt quality is as decisive as algorithm design in traditional planning',
    'Active tool-calling (Function Calling) enables spatially grounded reasoning',
    'Modular Skill YAML architecture → extensible to new tasks with zero code changes',
    'Semantic mapping emerges naturally as a side product of agent navigation',
], font_size=15)

# Right: Future Work
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
             'Future Directions', font_size=20, color=ACCENT_GREEN, bold=True)

add_bullet_slide_content(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5), [
    'Progressive Cognition: Scout → Inspector → Navigator multi-skill framework',
    'Persistent semantic memory across sessions (save/load JSON maps)',
    'Hybrid architecture: local lightweight model + cloud VLM to reduce latency',
    'Multi-robot formation: leader (perception) + followers (execution)',
    'Cross-domain generalization to outdoor / warehouse environments',
], font_size=15)

# Bottom: research interest alignment
add_rect(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0),
         RGBColor(0xFE, 0xF3, 0xE2))
add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11.5), Inches(0.8),
             'Research Interest Alignment:  This work sits at the intersection of Embodied AI, '
             'Vision-Language Models, and Multi-Robot Systems — areas where I am eager to '
             'pursue deeper research at the graduate level.',
             font_size=15, color=DARK_BLUE, bold=False)

add_page_number(slide, 9, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Slide 10: Thank You / Q&A
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_rect(slide, Inches(0), Inches(0), W, Inches(0.15), ACCENT_ORANGE)

add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2),
             'Thank You',
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.03), ACCENT_ORANGE)

add_text_box(slide, Inches(2.0), Inches(3.9), Inches(9.3), Inches(0.6),
             'Questions & Discussion',
             font_size=24, color=RGBColor(0xA0, 0xC8, 0xE8),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(1.5),
             '[Your Name]  |  [Email]  |  [University]\n'
             'GitHub: github.com/Kuko1414/FYP_Ros2',
             font_size=16, color=RGBColor(0x99, 0xBB, 0xDD),
             alignment=PP_ALIGN.CENTER, line_spacing=2.0)

add_page_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
output_path = '/home/kuko/humble_ws/Research_Experience_Presentation.pptx'
prs.save(output_path)
print(f'✅ PPT saved to: {output_path}')
